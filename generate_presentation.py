"""
Generador de Presentación PowerPoint para PBI Assistant
Reunión Hermes Agent - 2026-06-10
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Colores YPF (tema corporativo)
COLOR_YPF_YELLOW = RGBColor(242, 200, 17)
COLOR_YPF_BLACK = RGBColor(11, 14, 20)
COLOR_DARK_BG = RGBColor(17, 21, 32)
COLOR_TEXT_PRIMARY = RGBColor(232, 236, 244)
COLOR_TEXT_SECONDARY = RGBColor(139, 149, 168)
COLOR_ACCENT_BLUE = RGBColor(59, 130, 246)
COLOR_ACCENT_GREEN = RGBColor(16, 124, 16)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 1: PORTADA
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    shapes = slide.shapes

    # Background oscuro
    background = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_YPF_BLACK
    background.line.fill.background()

    # Título principal
    title_box = shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "PBI ASSISTANT"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLOR_YPF_YELLOW
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Chatbot de Power BI con IA para Usuarios No Técnicos"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    # Info inferior
    info_box = shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    info_frame = info_box.text_frame
    info_frame.text = "Integración con Hermes Agent • 2026-06-10"
    p = info_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 2: PROBLEMA
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "El Problema")

    add_bullet_list(slide, Inches(1), Inches(2), Inches(8), Inches(4), [
        "Usuarios de negocio necesitan consultar datos de Power BI",
        "No conocen DAX (lenguaje de consultas)",
        "No saben usar comandos técnicos de pbi-cli",
        "Desconocen la estructura del modelo semántico",
        "",
        "Resultado: Dependen de analistas técnicos → Cuellos de botella"
    ])

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 3: SOLUCIÓN
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "La Solución: PBI Assistant")

    add_bullet_list(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "Interfaz de chat con lenguaje natural",
        "Usuario escribe: 'Mostrame ventas por región'",
        "Claude CLI interpreta e invoca skill power-bi-dax",
        "pbi-cli ejecuta: EVALUATE SUMMARIZE(Sales, Region, [Total])",
        "Retorna tabla formateada con datos reales",
        "",
        "+ Feature Cards: 'Recetas' visuales que autocompletan prompts"
    ])

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 4: ARQUITECTURA - VISTA GENERAL
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Arquitectura de 3 Capas")

    # Diagrama simplificado con cajas
    left = Inches(2)
    width = Inches(6)
    height = Inches(0.8)

    # Capa 1: Frontend
    box1 = add_colored_box(slide, left, Inches(2), width, height, COLOR_ACCENT_BLUE)
    add_text_to_shape(box1, "1. FRONTEND (index.html - 1,941 líneas)\nHTML/CSS/JS puro • Feature cards • Markdown renderer", 14)

    # Flecha
    add_arrow(slide, Inches(5), Inches(2.9), Inches(5), Inches(3.3))
    add_small_text(slide, Inches(5.2), Inches(3), "HTTP REST")

    # Capa 2: Backend
    box2 = add_colored_box(slide, left, Inches(3.4), width, height, RGBColor(139, 92, 246))
    add_text_to_shape(box2, "2. BACKEND (server.py - 920 líneas)\nPython proxy • Auth • Logging • Session delegation", 14)

    # Flecha
    add_arrow(slide, Inches(5), Inches(4.3), Inches(5), Inches(4.7))
    add_small_text(slide, Inches(5.2), Inches(4.4), "subprocess")

    # Capa 3: IA
    box3 = add_colored_box(slide, left, Inches(4.8), width, height, COLOR_YPF_YELLOW)
    add_text_to_shape(box3, "3. IA + CLI (Claude CLI + pbi-cli)\nClaude: 13 skills • pbi-cli: Dual-layer (Model + Report)", 14, RGBColor(0, 0, 0))

    # Power BI
    add_small_text(slide, Inches(3.5), Inches(5.8), "↓\nPower BI Desktop / Archivos .pbip", center=True)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 5: FLUJO DE UN MENSAJE (8 PASOS)
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Flujo de un Mensaje (8 Pasos)")

    add_numbered_list(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(5), [
        "Usuario escribe → 'Mostrame ventas por región'",
        "Frontend → POST /api/chat {message, tone, session_id}",
        "Backend → Escribe mensaje a temp file",
        "Backend → subprocess: claude --resume {session_id} < msg.txt",
        "Claude CLI → Invoca skill power-bi-dax",
        "pbi-cli → Ejecuta: pbi --json dax execute -e 'EVALUATE...'",
        "Claude CLI → Formatea respuesta en español (según tone)",
        "Frontend → Renderiza Markdown como tabla HTML"
    ], 16)

    add_small_text(slide, Inches(7), Inches(6.5), "⏱️ Tiempo total: 2-5 segundos", center=True)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 6: 5 INNOVACIONES CLAVE
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "5 Innovaciones Clave")

    innovations = [
        ("1️⃣ Session Persistence", "--resume {session_id} → Claude recuerda todo\nFrontend no reenvía historial"),
        ("2️⃣ Feature Cards", "20 recetas predefinidas, 5 categorías\nClick → Autocompleta → Envía"),
        ("3️⃣ Auto-Detection", "Detecta: Power BI abierto, múltiples conexiones, .pbip files\nUsuario no configura nada"),
        ("4️⃣ Adaptive Tone", "Porteño (voseo), Formal (usted), Neutral\nSe adapta a cultura organizacional"),
        ("5️⃣ Usage Metrics", "Logs JSONL + Dashboard admin + Export CSV\nCompliance, insights, power users")
    ]

    y_pos = 2
    for title, desc in innovations:
        box = add_colored_box(slide, Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.75), COLOR_DARK_BG)
        add_text_to_shape(box, f"{title}\n{desc}", 12)
        y_pos += 0.95

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 7: CÓMO LEE Y ESCRIBE ARCHIVOS
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Cómo Lee y Escribe Archivos")

    # Modelo Semántico
    box1 = add_colored_box(slide, Inches(0.8), Inches(2), Inches(4), Inches(2.2), COLOR_ACCENT_BLUE)
    add_text_to_shape(box1, "MODELO SEMÁNTICO\n(Requiere Power BI Desktop)\n\n• TCP → Analysis Services\n• En memoria (tiempo real)\n• DAX queries\n• Metadata operations", 14)

    # Reportes
    box2 = add_colored_box(slide, Inches(5.2), Inches(2), Inches(4), Inches(2.2), COLOR_ACCENT_GREEN)
    add_text_to_shape(box2, "CAPA DE REPORTES\n(Trabaja offline)\n\n• File I/O → .pbip JSON\n• Offline (no Power BI)\n• Visuales, páginas, temas\n• Modifica report.json", 14)

    # Formato .pbip
    add_small_text(slide, Inches(2), Inches(4.5),
        "MiReporte.pbip\n├── *.SemanticModel/\n│   └── definition/\n└── *.Report/\n    └── definition/report.json",
        monospace=True)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 8: VIBE CODE
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Construido con Vibe Code")

    add_bullet_list(slide, Inches(1), Inches(2), Inches(8), Inches(2), [
        "Vibe Code = Desarrollo guiado por IA",
        "No escribimos código línea por línea",
        "Claude generó la aplicación completa en ~3 días"
    ])

    # Proceso
    box = add_colored_box(slide, Inches(1.5), Inches(4.5), Inches(7), Inches(2.5), COLOR_DARK_BG)
    add_text_to_shape(box,
        "PROCESO:\n"
        "1. Mockup inicial → Dibujamos UI\n"
        "2. Prompt → 'Crea chatbot de Power BI con...'\n"
        "3. Claude genera → index.html (1,941 líneas)\n"
        "4. Iteramos → 'Agrega selector de tono'\n"
        "5. Claude refina → Mantiene consistencia", 13)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 9: MÉTRICAS DE USO REAL
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Métricas de Uso Real (2 Semanas)")

    # KPIs en grid
    kpis = [
        ("450+", "Mensajes", COLOR_YPF_YELLOW),
        ("23", "Sesiones", COLOR_ACCENT_BLUE),
        ("8", "Usuarios", COLOR_ACCENT_GREEN),
        ("3.5s", "Latencia Avg", RGBColor(236, 72, 153))
    ]

    x_pos = 1.2
    for value, label, color in kpis:
        box = add_colored_box(slide, Inches(x_pos), Inches(2.2), Inches(1.8), Inches(1.2), color)
        add_text_to_shape(box, f"{value}\n{label}", 16, bold_first_line=True)
        x_pos += 2

    # Insights
    add_bullet_list(slide, Inches(1), Inches(4), Inches(8), Inches(2.5), [
        "Top skills: dax execute (40%), measure list (25%), visual add (15%)",
        "Power users: 2 usuarios = 60% del tráfico",
        "Horarios pico: 10-12am, 3-5pm"
    ])

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 10: SEGURIDAD
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Seguridad y Compliance")

    security_items = [
        ("🔐 Whitelist de Usuarios", "config.json → Solo usuarios autorizados"),
        ("🛡️ Prompt Injection Protection", "Restricción: Solo consultas de Power BI"),
        ("📝 Audit Trail", "Logs JSONL: timestamp, usuario, mensaje, duración"),
        ("🔒 Row-Level Security", "Heredada de Power BI (RLS)"),
        ("🏠 On-Premise Deployment", "Claude CLI local (no cloud)")
    ]

    y_pos = 2.2
    for title, desc in security_items:
        box = add_colored_box(slide, Inches(1.5), Inches(y_pos), Inches(7), Inches(0.7), COLOR_DARK_BG)
        add_text_to_shape(box, f"{title}\n{desc}", 12)
        y_pos += 0.9

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 11: MIGRACIÓN A HERMES - OPCIÓN 1
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Integración Hermes: Opción 1 - Coexistencia")

    # Diagrama
    box1 = add_colored_box(slide, Inches(2), Inches(2), Inches(6), Inches(0.8), COLOR_ACCENT_BLUE)
    add_text_to_shape(box1, "Hermes Agent (Orchestrator)", 14)

    add_arrow(slide, Inches(5), Inches(2.9), Inches(5), Inches(3.3))
    add_small_text(slide, Inches(5.2), Inches(3), "HTTP call")

    box2 = add_colored_box(slide, Inches(2), Inches(3.4), Inches(6), Inches(0.8), COLOR_YPF_YELLOW)
    add_text_to_shape(box2, "PBI Assistant (actual, vía /api/chat)", 14, RGBColor(0, 0, 0))

    # Pros/Cons
    add_pros_cons(slide, Inches(1), Inches(4.5),
        ["Sin reescribir código", "Multi-dominio", "PBI sigue standalone"],
        ["Doble capa IA", "Latencia HTTP"])

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 12: MIGRACIÓN A HERMES - OPCIÓN 2
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Integración Hermes: Opción 2 - Migración Completa")

    # Diagrama
    boxes = [
        ("Hermes Agent SDK", COLOR_ACCENT_BLUE, 2),
        ("anthropic.Client()", RGBColor(139, 92, 246), 2.9),
        ("Custom tools (pbi-cli wrappers)", COLOR_ACCENT_GREEN, 3.8),
        ("pbi-cli → Power BI", COLOR_YPF_YELLOW, 4.7)
    ]

    for text, color, y in boxes:
        box = add_colored_box(slide, Inches(2.5), Inches(y), Inches(5), Inches(0.7), color)
        text_color = RGBColor(0, 0, 0) if color == COLOR_YPF_YELLOW else COLOR_TEXT_PRIMARY
        add_text_to_shape(box, text, 13, text_color)
        if y < 4.7:
            add_arrow(slide, Inches(5), Inches(y+0.75), Inches(5), Inches(y+0.95))

    # Pros/Cons
    add_pros_cons(slide, Inches(0.8), Inches(5.8),
        ["Control total", "Una capa IA", "Métricas nativas"],
        ["Más código", "Session manual"])

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 13: MIGRACIÓN A HERMES - OPCIÓN 3
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Integración Hermes: Opción 3 - Híbrido")

    add_bullet_list(slide, Inches(1), Inches(2), Inches(8), Inches(2.5), [
        "Lo mejor de ambos mundos",
        "Hermes Agent para orquestación multi-dominio",
        "pbi-cli tools convertidos a Hermes @tool decorators",
        "Mantiene arquitectura modular",
        "Migración incremental sin riesgo"
    ])

    # Ejemplo de código
    box = add_colored_box(slide, Inches(1.5), Inches(4.7), Inches(7), Inches(2), COLOR_DARK_BG)
    add_text_to_shape(box,
        "@tool\n"
        "def execute_dax_query(query: str) -> dict:\n"
        "    '''Execute DAX query against Power BI'''\n"
        "    result = subprocess.run(['pbi', 'dax', 'execute', '-e', query])\n"
        "    return json.loads(result.stdout)", 11, monospace=True)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 14: COSTOS Y ROI
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Costos y ROI")

    # Costos
    box1 = add_colored_box(slide, Inches(1), Inches(2), Inches(4), Inches(2.5), COLOR_DARK_BG)
    add_text_to_shape(box1,
        "COSTOS MENSUALES\n(50 usuarios activos)\n\n"
        "• Claude API: $45\n"
        "• Compute (2CPU+4GB): $50\n"
        "• Storage: $5\n"
        "────────────\n"
        "TOTAL: $100-200/mes", 13)

    # ROI
    box2 = add_colored_box(slide, Inches(5.2), Inches(2), Inches(3.8), Inches(2.5), COLOR_ACCENT_GREEN)
    add_text_to_shape(box2,
        "ROI ESTIMADO\n\n"
        "Ahorro: 12.5 hrs/semana\n"
        "Valor: $2,500/mes\n\n"
        "ROI = 1,150%\n"
        "🚀", 14, bold_first_line=True)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 15: PRÓXIMOS PASOS
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Próximos Pasos")

    add_numbered_list(slide, Inches(1.5), Inches(2.2), Inches(7), Inches(4.5), [
        "DECISIÓN: ¿Coexistencia, migración o híbrido?",
        "PoC: Integrar 1 skill de PBI como Hermes tool",
        "COMPARACIÓN: Latencia, developer experience, mantenibilidad",
        "ROLLOUT: Plan de migración si decidimos full Hermes",
        "ESCALAR: Multi-tenancy, rate limiting, load balancing"
    ], 18)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 16: COMPARACIÓN TÉCNICA
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Comparación Técnica")

    # Tabla comparativa
    table_data = [
        ("CARACTERÍSTICA", "CLAUDE CLI (Actual)", "HERMES (Migrado)"),
        ("Session Management", "Automático (--resume)", "Manual (Dict)"),
        ("Tools Definition", "Skills predefinidos", "Define cada tool"),
        ("Debugging", "stdout legible", "JSON responses"),
        ("Latency", "subprocess overhead", "Direct HTTP"),
        ("Dependency", "Claude CLI binary", "pip package"),
        ("Multi-tenancy", "Difícil (1 CLI)", "Fácil (N clients)"),
        ("Cloud Deploy", "Requiere instalación", "Solo pip install")
    ]

    add_comparison_table(slide, Inches(0.8), Inches(2), Inches(8.4), table_data)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 17: CONCLUSIÓN
    # ═══════════════════════════════════════════════════════════════════════
    slide = add_title_slide(prs, "Conclusión")

    conclusion_box = add_colored_box(slide, Inches(1), Inches(2), Inches(8), Inches(4), COLOR_YPF_YELLOW)
    add_text_to_shape(conclusion_box,
        "PBI Assistant demuestra que pbi-cli + Claude CLI\n"
        "es una combinación poderosa para democratizar Power BI.\n\n"
        "✅ Construido en 3 días con Vibe Code\n"
        "✅ 8 usuarios activos, 450+ mensajes procesados\n"
        "✅ ROI estimado de 1,150%\n\n"
        "Con Hermes Agent, podemos evolucionar de un asistente\n"
        "single-purpose a un orquestador multi-dominio.\n\n"
        "La arquitectura modular permite migración incremental sin riesgo.",
        16, RGBColor(0, 0, 0))

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 18: Q&A
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    # Background
    background = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_YPF_BLACK
    background.line.fill.background()

    # Q&A texto
    qa_box = shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    qa_frame = qa_box.text_frame
    qa_frame.text = "¿Preguntas?"
    p = qa_frame.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLOR_YPF_YELLOW
    p.alignment = PP_ALIGN.CENTER

    # Info
    info_box = shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
    info_frame = info_box.text_frame
    info_frame.text = "Documentación completa disponible en:\npbi-cli/web/README_REUNION_HERMES.md"
    for p in info_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT_SECONDARY
        p.alignment = PP_ALIGN.CENTER

    return prs


# ═══════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════

def add_title_slide(prs, title_text):
    """Add slide with title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    # Background
    background = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_YPF_BLACK
    background.line.fill.background()

    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_YPF_YELLOW

    # Underline
    line = shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_YPF_YELLOW
    line.line.fill.background()

    return slide


def add_bullet_list(slide, left, top, width, height, items, font_size=16):
    """Add bullet list to slide."""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        p.level = 0

        if item.strip():  # Solo bullet si no es línea vacía
            p.bullet = True


def add_numbered_list(slide, left, top, width, height, items, font_size=16):
    """Add numbered list to slide."""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        p.level = 0


def add_colored_box(slide, left, top, width, height, color):
    """Add colored rectangle."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    return box


def add_text_to_shape(shape, text, font_size=14, color=COLOR_TEXT_PRIMARY, bold_first_line=False, monospace=False):
    """Add text to a shape."""
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        if monospace:
            p.font.name = 'Consolas'

        if bold_first_line and i == 0:
            p.font.bold = True


def add_small_text(slide, left, top, text, font_size=12, center=False, monospace=False):
    """Add small text box."""
    box = slide.shapes.add_textbox(left, top, Inches(3), Inches(1))
    text_frame = box.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    if center:
        p.alignment = PP_ALIGN.CENTER
    if monospace:
        p.font.name = 'Consolas'


def add_arrow(slide, x1, y1, x2, y2):
    """Add arrow connector."""
    from pptx.enum.shapes import MSO_CONNECTOR
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = COLOR_TEXT_SECONDARY
    connector.line.width = Pt(2)


def add_pros_cons(slide, left, top, pros, cons):
    """Add pros/cons boxes."""
    # Pros
    pros_box = add_colored_box(slide, left, top, Inches(4), Inches(1.8), RGBColor(20, 60, 20))
    text_frame = pros_box.text_frame
    text_frame.text = "✅ PROS:\n" + "\n".join(f"• {p}" for p in pros)
    for p in text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # Cons
    cons_box = add_colored_box(slide, left + Inches(4.5), top, Inches(4), Inches(1.8), RGBColor(60, 20, 20))
    text_frame = cons_box.text_frame
    text_frame.text = "❌ CONS:\n" + "\n".join(f"• {c}" for c in cons)
    for p in text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_PRIMARY


def add_comparison_table(slide, left, top, width, data):
    """Add comparison table."""
    from pptx.util import Inches

    rows = len(data)
    cols = len(data[0])

    table = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.4 * rows)).table

    # Header row
    for col_idx, cell_text in enumerate(data[0]):
        cell = table.rows[0].cells[col_idx]
        cell.text = cell_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_YPF_YELLOW

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx in range(1, rows):
        for col_idx, cell_text in enumerate(data[row_idx]):
            cell = table.rows[row_idx].cells[col_idx]
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_DARK_BG

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = COLOR_TEXT_PRIMARY
                if col_idx == 0:
                    paragraph.font.bold = True


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import sys
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    print("Generando presentación PowerPoint...")
    prs = create_presentation()

    output_file = "PBI_Assistant_Reunion_Hermes.pptx"
    prs.save(output_file)

    print(f"✅ Presentación creada: {output_file}")
    print(f"📊 Total de slides: {len(prs.slides)}")
    print("\nSlides generadas:")
    print("  1. Portada")
    print("  2. El Problema")
    print("  3. La Solución")
    print("  4. Arquitectura de 3 Capas")
    print("  5. Flujo de un Mensaje")
    print("  6. 5 Innovaciones Clave")
    print("  7. Cómo Lee y Escribe Archivos")
    print("  8. Construido con Vibe Code")
    print("  9. Métricas de Uso Real")
    print(" 10. Seguridad y Compliance")
    print(" 11. Integración Hermes: Opción 1")
    print(" 12. Integración Hermes: Opción 2")
    print(" 13. Integración Hermes: Opción 3")
    print(" 14. Costos y ROI")
    print(" 15. Próximos Pasos")
    print(" 16. Comparación Técnica")
    print(" 17. Conclusión")
    print(" 18. Q&A")
